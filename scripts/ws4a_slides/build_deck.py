#!/usr/bin/env python3
"""Build the WS4A slide deck (.pptx).

Each slide: a SAY-THIS headline on top, the figure in the middle, a plain-language
line at the bottom, and the full speaking script + Q&A in the speaker notes.

    python build_deck.py <pack_dir> <out.pptx>
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

PACK = Path(sys.argv[1])
OUT = Path(sys.argv[2])
FIG = PACK / "00_THE_THREE_SLIDES"

W, H = Inches(13.333), Inches(7.5)
INK   = RGBColor(0x1A, 0x1A, 0x1A)
MUT   = RGBColor(0x5A, 0x5A, 0x58)
BG    = RGBColor(0xFC, 0xFC, 0xFB)
ACC   = RGBColor(0x2A, 0x78, 0xD6)
RED   = RGBColor(0xB0, 0x24, 0x18)


def textbox(slide, x, y, w, h, text, size, color=INK, bold=False,
            align=PP_ALIGN.LEFT, italic=False, space_after=4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def picture_fit(slide, img, x, y, max_w, max_h):
    """Insert centred in the box, preserving aspect ratio."""
    iw, ih = Image.open(img).size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(img), int(x + (max_w - w) / 2),
                             int(y + (max_h - h) / 2), width=w, height=h)


def blocks_box(slide, x, y, w, h, blocks, size=11.5):
    """Labelled paragraphs: LABEL in accent colour, then the sentence in ink.

    The deck is presented by someone who did not run the analysis, so the slide has to
    carry the explanation itself -- the speaker notes are a bonus, not the mechanism.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for label, body in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7)
        r = p.add_run(); r.text = f"{label}  "
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = ACC
        r.font.name = "Calibri"
        r2 = p.add_run(); r2.text = body
        r2.font.size = Pt(size); r2.font.color.rgb = INK; r2.font.name = "Calibri"
    return tb


def add_slide(prs, *, kicker, say, figure, blocks, notes, subtitle=None,
              fig_h=3.0, text_size=11.5):
    """Kicker, a SHORT conclusion as the title, an elaborating line, then the figure
    and the labelled explanation blocks.

    The title states what the slide concludes -- not what it contains -- so a reader
    who sees only the titles still gets the argument. The nuance moves to `subtitle`.

    The figure is deliberately given less room than the text: the presenter can resize
    it by hand, but they cannot invent the explanation.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    if kicker:
        textbox(s, Inches(0.5), Inches(0.14), Inches(12.4), Inches(0.3),
                kicker, 12, ACC, bold=True)
    textbox(s, Inches(0.5), Inches(0.40), Inches(12.4), Inches(0.55),
            say, 26, INK, bold=True)
    top = 1.12
    if subtitle:
        textbox(s, Inches(0.5), Inches(0.98), Inches(12.4), Inches(0.35),
                subtitle, 14, MUT)
        top = 1.48
    if figure:
        picture_fit(s, figure, Inches(0.45), Inches(top),
                    Emu(int(Inches(12.45))), Emu(int(Inches(fig_h))))
    blocks_box(s, Inches(0.5), Inches(top + fig_h + 0.18), Inches(12.4),
               Inches(7.35 - (top + fig_h + 0.18)), blocks, text_size)
    s.notes_slide.notes_text_frame.text = notes
    return s


prs = Presentation()
prs.slide_width, prs.slide_height = W, H

# ---------------------------------------------------------------- TITLE
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
textbox(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(1.3),
        "Cell Painting tells you about mechanism —\nnot about toxicity", 34, INK,
        bold=True)
textbox(s, Inches(0.9), Inches(3.45), Inches(11.5), Inches(0.5),
        "…and it beats chemical structure, which is free, at doing so.", 17, ACC,
        bold=True)
textbox(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(1.1),
        "WS4A · cross-modal integration · HepG2\n"
        "119 compounds · morphology (636) × gene expression (41,780) × "
        "chemical structure (1,024)", 15, MUT)
textbox(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(1.6),
        "SAY THIS FIRST, OR NOTHING ELSE MAKES SENSE:\n"
        "Every number in this section is an HONEST GAP — how well the model did on the "
        "real answers, minus how well the same model did when we scrambled the answers "
        "between compounds. Scrambled means there is nothing left to learn, so that "
        "second number is what luck alone buys.  Across our grid, luck scored between "
        "0.38 and 0.60 — where 0.50 is a coin flip.  That is why we never quote a raw "
        "score.", 13, RED, bold=False)
s.notes_slide.notes_text_frame.text = (
    "OPENING LINE — say this before anything else, or no axis on any later slide "
    "means anything:\n\n"
    "\"Every number I show is an HONEST GAP: how well the model did on the real "
    "answers, minus how well the identical model did when we scrambled the answers "
    "between compounds. Scrambled means there is nothing left to learn, so that "
    "second number is what luck alone buys. Across our whole grid, luck scored "
    "between 0.38 and 0.60 — where 0.50 is a coin flip. That is why we never quote a "
    "raw score.\"\n\n"
    "THE THREE NUMBERS TO MEMORISE\n"
    "  +0.143       what morphology adds to chemistry for mechanism (4/4 models)\n"
    "  -0.273 vs +0.017   cost of dropping to 2,000 HVGs vs what the leakage was worth\n"
    "  0 of 41,780  genes surviving the shared-axis permutation null")

# ---------------------------------------------------------------- METHOD
add_slide(
    prs,
    kicker="WHAT WAS TRAINED, AND HOW",
    say="Only 119 drugs — so every score gets a scrambled-label twin",
    subtitle="Six questions · five kinds of data · four model families · every setting "
             "chosen on drugs the score never sees.",
    figure=FIG / "SLIDE_0_how_it_was_trained.png",
    fig_h=2.55,
    blocks=[
        ("THE DATA.",
         "119 compounds, each described three ways: its chemical structure (free — "
         "computed from the molecule, no experiment); the Cell Painting photograph of "
         "treated HepG2 cells (636 measurements of shape, texture and where organelles "
         "sit); and gene expression (41,780 genes). The expression comes from 384,533 "
         "single cells, but cells given the same drug are replicates — they are "
         "averaged into ONE row per drug. So the model sees 119 examples, not 384,000."),
        ("WHAT WE PREDICT.",
         "Six questions: is this a DNA-synthesis inhibitor (14 yes vs 50 no), and is it "
         "toxic to heart, lung, kidney, liver or fertility (~70 drugs each). Two further "
         "toxicity endpoints were refused outright — at 68 yes vs 2 no, a model that "
         "always says yes is right 97% of the time and has learned nothing."),
        ("HOW A MODEL IS TESTED.",
         "Hide a fifth of the drugs, learn from the rest, guess the hidden ones, and "
         "rotate which are hidden 25 times. Every setting — how much to penalise, how "
         "many features to keep, how deep a tree may grow — is chosen using only the "
         "training drugs, so the ones being scored never influence anything."),
        ("WHY EVERYTHING RUNS TWICE.",
         "With 119 rows and up to 41,780 measurements, some measurements line up with "
         "the answer by pure chance. So we repeat the whole thing with the answers "
         "SCRAMBLED between drugs — nothing left to learn. Luck alone scored 0.38 to "
         "0.60, where 0.50 is a coin flip. Every number on the next slides is real "
         "minus scrambled."),
    ],
    notes=(
        "SAY:\n"
        "\"First, exactly what we did.\n\n"
        "Panel 1 — the questions. Mechanism of action: is this a DNA-synthesis "
        "inhibitor, 14 yes against 50 no. And five toxicity endpoints, yes/no, about "
        "70 drugs each. Two further endpoints we REFUSED to score at all — at 68 to 2, "
        "a model that always says yes is right 97 percent of the time and has learned "
        "nothing.\n\n"
        "Panel 2 — what the model may look at. Chemistry, the photograph, gene "
        "expression, and two combinations. Chemistry is our control because it is "
        "free: you compute it from the molecule, no experiment at all.\n\n"
        "Panel 3 — four model families. Three are variations on drawing a straight "
        "dividing line; XGBoost builds decision trees, a completely different way of "
        "looking. That matters: when methods that disagree about HOW to find structure "
        "agree about WHERE it is, that means something.\n\n"
        "Panel 4 — this is the part people usually get wrong. The orange blocks are "
        "the drugs we score on. They never influence anything that gets fitted. Every "
        "setting — how much to penalise, how many features to keep, how deep a tree may "
        "grow — is chosen by splitting the BLUE part again and looking only in there.\n\n"
        "Panel 5 — we did it twice: once with a fixed grid declared in advance, once "
        "letting Optuna search 20 combinations per fold. Both kept, both compared. And "
        "critically, the scrambled control gets the same 20 trials — if you tune the "
        "real answers harder than the control, you invent a result.\n\n"
        "Panel 6 — then the whole thing again with the answers scrambled between drugs. "
        "Nothing left to learn, so whatever it scores is luck. On our grid luck alone "
        "scored 0.38 to 0.60, where 0.5 is a coin flip. So we never quote a score — we "
        "quote real minus scrambled.\"\n\n"
        "IF ASKED\n"
        "Q: Why 119 rows when you have 384,000 cells?\n"
        "A: Cells given the same drug are replicates, not independent observations. "
        "Treating them as 384,000 examples would inflate everything enormously.\n\n"
        "Q: How were hyperparameters chosen?\n"
        "A: Entirely inside the inner loop, on the training part of each outer fold. "
        "The scored drugs never influenced them. A pre-flight check asserts this by "
        "tagging every row and verifying no search ever saw all of them.\n\n"
        "Q: Did tuning just make the numbers look better?\n"
        "A: We measured that. The shuffled-label scores moved by −0.0003, and zero of "
        "120 comparisons were flagged as bias. It also produced no new findings — 56 "
        "results cleared zero before tuning and 56 after.\n\n"
        "Q: Isn't this just a permutation test?\n"
        "A: Yes, but run separately for every model × question × data-type combination "
        "and subtracted individually, because each has its own luck level."))

# ---------------------------------------------------------------- RESULT
add_slide(
    prs,
    kicker="RESULT 1 OF 3",
    say="Morphology beats chemistry — for mechanism only",
    subtitle="+0.143 over chemical structure for mechanism of action, all four models "
             "agreeing.  For all five toxicity endpoints it adds nothing.",
    figure=FIG / "SLIDE_1_ALT_signal_map.png",
    fig_h=2.85,
    blocks=[
        ("HOW TO READ IT.",
         "Rows are the six questions; columns are what the model was allowed to look "
         "at. Each cell holds the best of four models, measured as the HONEST GAP — its "
         "score on the real answers minus its score when the answers were scrambled. "
         "Red = real signal. White = nothing. The small text under each number is how "
         "many of the four models independently found it."),
        ("WHAT IT SHOWS.",
         "The top row — mechanism of action — is the only row that works across the "
         "board. Chemical structure alone gives +0.175; the photograph alone gives "
         "+0.267; the two together give +0.318. Now read DOWN the morphology column: "
         "every toxicity row is near-white (+0.012, +0.011, +0.008, +0.064, −0.005) "
         "with 0 of 4 models finding anything. Meanwhile the chemistry column IS "
         "predicting the toxicity endpoints."),
        ("THE POINT.",
         "Chemistry is the control because it is free — no cells, no microscope, no "
         "sequencing — and chemists have predicted drug properties from structure for "
         "thirty years. So the question is not whether imaging beats chance, but "
         "whether it beats the free option. For MECHANISM it does, adding +0.143 on top "
         "of chemistry with all four models agreeing. For TOXICITY it does not."),
    ],
    notes=(
        "SAY:\n"
        "\"Chemistry is our control, and it is the right control because it is FREE — "
        "you compute a structural fingerprint from the molecule on a laptop, no cells, "
        "no microscope, no sequencing. Chemists have predicted drug properties this way "
        "for thirty years. So 'imaging beats random guessing' is not a finding anyone "
        "should act on. The question is whether imaging beats the free option.\n\n"
        "Top row is mechanism of action. Chemistry alone: 0.175. The photograph alone: "
        "0.267 — better. Both together: 0.318. So the photograph adds 0.143 on top of "
        "structure, and the small print says all four of our models found it "
        "independently.\n\n"
        "Now look down the morphology column at the five toxicity rows: 0.012, 0.011, "
        "0.008, 0.064, minus 0.005. Nothing, and zero of four models in every case. "
        "Meanwhile the chemistry column IS predicting toxicity.\n\n"
        "So the claim is narrow and defensible: Cell Painting tells you about "
        "MECHANISM, over and above chemical structure. It does not tell you about "
        "toxicity.\"\n\n"
        "IF ASKED\n"
        "Q: Isn't 0.143 small?\n"
        "A: It is an 82% increase over the chemistry-only gap, and four independent "
        "model families found it while the scrambled control did not.\n\n"
        "Q: Why does ADDING data make toxicity worse?\n"
        "A: Small-sample dilution. Adding 636 columns to ~70 compounds gives the model "
        "636 more chances to fit a coincidence.\n\n"
        "Q: So Cell Painting can't see toxicity?\n"
        "A: Careful — say 'not from 70 compounds with these labels'. 'Known to be "
        "hepatotoxic' is a coarse literature annotation, much noisier than a curated "
        "mechanism class.\n\n"
        "DO NOT SAY: anything about A549 (its morphology is contaminated), and never "
        "quote a raw score without its scrambled twin."))

# ---------------------------------------------------------------- HVG
add_slide(
    prs,
    kicker="RESULT 2 OF 3  ·  reconciles with Task 1",
    say="Gene selection, not the modality, is why expression lost",
    subtitle="Dropping to 2,000 highly variable genes costs 16× what the leakage it is "
             "usually criticised for was worth.",
    figure=FIG / "SLIDE_2_hvg_selection_cost.png",
    fig_h=2.75,
    blocks=[
        ("THE BACKGROUND.",
         "Almost every single-cell workflow keeps only the ~2,000 most variable genes "
         "out of 40,000. Task 1 of this workstream does that, and its own caveat slide "
         "flags the leakage. But that caveat bundles TWO different problems: REDUCTION "
         "(throwing 39,780 genes away) and LEAKAGE (choosing which 2,000 using drugs "
         "you later test on). Nobody knew which was actually costing anything."),
        ("HOW TO READ IT.",
         "Left: three groups of bars, identical cross-validation in each — all 41,780 "
         "genes; 2,000 genes chosen separately INSIDE each training fold (leakage-free); "
         "and 2,000 chosen once using all drugs (the bundled approach). The brown dotted "
         "line is the best morphology result. Right: the two effects separated — blue is "
         "what reduction costs, red is what the leakage was worth."),
        ("THE POINT.",
         "Reduction costs −0.273; the leakage was worth +0.017. About sixteen times "
         "more. And with all 41,780 genes the expression arm clears the brown line — "
         "with 2,000 it does not come close. So \"gene expression does not beat "
         "morphology\" looks like a property of the gene selection, not of the modality. "
         "Both pipelines agree once the same genes are used."),
    ],
    notes=(
        "LEAD WITH THE AGREEMENT — this reconciles two analyses, it does not overturn "
        "one.\n\n"
        "SAY:\n"
        "\"First, where we agree with the main analysis. It finds that fusing "
        "morphology and gene expression does not improve mechanism prediction. We got "
        "the same thing independently — our fusion gain is plus 0.005, theirs is minus "
        "0.096. Different pipeline, different task formulation, same conclusion. That "
        "is a replication and it is worth saying out loud.\n\n"
        "Second, one thing we can add. Task 1 selects 2,000 highly variable genes "
        "before cross-validating, and flags the leakage in its own caveat. But that "
        "caveat bundles TWO different problems. One is leakage — choosing which genes "
        "using compounds you later test on. The other is reduction — simply throwing "
        "39,780 genes away. Nobody knew which was costing anything, so we separated "
        "them.\n\n"
        "Blue is reduction. Red is leakage. Reduction costs about sixteen times more.\n\n"
        "And the brown dotted line is the best morphology arm. With all 41,780 genes, "
        "expression clears it. With 2,000 highly variable genes it does not come close. "
        "So the finding that gene expression does not beat morphology looks like a "
        "property of the gene selection, not of the modality. Our two pipelines agree "
        "once the same genes are used.\"\n\n"
        "CAVEAT TO STATE: this ran on mechanism only, one target, untuned — which is "
        "why the morphology reference is 0.233 and not 0.267. Everything in the figure "
        "is like-for-like.\n\n"
        "BACKGROUND IF PUSHED: Ambroise & McLachlan (PNAS 2002) showed gene selection "
        "before cross-validation gives near-zero error rates on data with NO signal. "
        "That is exactly the leakage the caveat flags — and we are showing it was the "
        "smaller of the two problems.\n\n"
        "The x marks XGBoost, which was degenerate in every arm (predicted one class "
        "every fold) and is therefore not a measurement."))

# ---------------------------------------------------------------- CONCORDANCE
add_slide(
    prs,
    kicker="RESULT 3 OF 3  ·  why this is not one lucky model",
    say="Four different models, the same answer",
    subtitle="Model families that fail in different ways agree on where the signal is — "
             "better evidence here than a p-value cross-validation cannot honestly give.",
    figure=FIG / "SLIDE_3_models_agree.png",
    fig_h=2.65,
    blocks=[
        ("WHY THIS SLIDE.",
         "The obvious objection to any result from 119 compounds is that we got lucky "
         "with one model. Our uncertainty intervals are also anti-conservative — "
         "cross-validation folds share training data, so there is no honest variance "
         "estimate at this sample size. So instead of a fragile p-value we show four "
         "model families that fail in different ways, and ask whether they agree."),
        ("HOW TO READ IT.",
         "Left: how similarly each pair of models ranks all 30 question × data "
         "combinations (1.00 = identical ordering, 0 = unrelated). Right: every "
         "combination as one row with all four models plotted, sorted by its WEAKEST "
         "model. The red line is zero — no better than that model's own scrambled "
         "control. Green and shaded = all four models clear zero."),
        ("THE POINT.",
         "The three linear models agree closely (0.73–0.87). XGBoost, which builds "
         "decision trees rather than drawing a line, agrees far less (0.47–0.53) — and "
         "that is exactly why its agreement counts: it is looking for a different kind "
         "of pattern. In 14 of the 30 combinations all four models clear zero, and that "
         "green set includes mechanism from morphology. Most of the other 16 sit at "
         "zero: the signal is concentrated, not everywhere."),
    ],
    notes=(
        "SAY:\n"
        "\"The obvious objection to anything from 119 compounds is that we got lucky "
        "with one model. So we ran four families that look for structure in genuinely "
        "different ways.\n\n"
        "On the left, how similarly each pair ranks all thirty combinations. One point "
        "zero would be identical ordering. The three linear models — elastic net, "
        "linear SVM, sparse PLS-DA — agree at 0.73 to 0.87. They are all variations on "
        "drawing a dividing line, so that much agreement is expected.\n\n"
        "XGBoost agrees much less, around 0.5. It builds decision trees. And that low "
        "number is the interesting one: if all four sat at 0.95 you would have four "
        "views of essentially one model, and 'they all agree' would mean nothing. At "
        "0.5, when it does land on the same combinations as the others, that is "
        "independent confirmation.\n\n"
        "On the right, all thirty combinations, sorted by their WEAKEST model — a "
        "combination is only as good as the worst of the four. Most of them sit at "
        "zero: the signal is concentrated at the top, which is what an honest landscape "
        "looks like. In the shaded band, 14 of 30, every single model clears zero — and "
        "that band includes mechanism of action from morphology, our headline.\"\n\n"
        "WHY THIS SLIDE, AND WHY IT MATTERS MORE THAN A P-VALUE HERE\n"
        "Our uncertainty intervals are anti-conservative — cross-validation folds share "
        "training data, so there is no unbiased variance estimate for k-fold CV "
        "(Bengio & Grandvalet 2004). Rather than lean on a shaky interval, we lean on "
        "four methods with different failure modes agreeing. At this sample size that "
        "is the more honest evidence.\n\n"
        "BE HONEST ABOUT THE ALTERNATIVE READING IF PUSHED\n"
        "A rho of 0.5 for XGBoost can mean two things: it is finding different REAL "
        "structure, or it is simply noisier and worse on this data. Both are consistent "
        "with 0.5 — XGBoost had 10 degenerate rows in the untuned run and rarely wins a "
        "combination. So the safe claim is the narrow one: the three linear models "
        "agree strongly, and on the top combinations XGBoost independently clears zero "
        "too. Do not oversell the 'different inductive bias' story.\n\n"
        "WHY SORTED BY THE WEAKEST MODEL: sorting by the mean would scatter the "
        "'all four clear zero' rows through the panel and the shading would read as "
        "stripes. Sorting by the minimum makes that set exactly the top block, and it "
        "is the conservative ranking anyway.\n\n"
        "IF ASKED WHY ONLY SIX ROWS ARE LABELLED: 30 labels do not read from the back "
        "of an auditorium -- the panel is dense, not merely small, so a bigger screen "
        "enlarges the rows AND the gaps and never improves the ratio. All 30 rows are "
        "still plotted, because you do not need to read a label to see the trend. The "
        "fully labelled version is in the pack at 04_methods/08_model_concordance.png.\n\n"
        "IF ASKED\n"
        "Q: Why not just report a p-value?\n"
        "A: See above — the interval for a k-fold CV mean is not trustworthy at this n. "
        "Model consensus is the substitute, and we say so.\n\n"
        "Q: What happened to the gene-level analysis?\n"
        "A: We have it — a permutation null on the shared morphology-expression axis "
        "where 0 of 41,780 genes survive. We left it out because this deck already "
        "reports 0 of 1.27 million feature-gene pairs at FDR below 0.05, by a different "
        "method. Same conclusion, already covered. It is in the backup."))

# ---------------------------------------------------------------- CLOSE
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
textbox(s, Inches(0.6), Inches(0.42), Inches(12.2), Inches(0.55),
        "Narrow claims, and the controls that made them narrow", 26, INK, bold=True)
textbox(s, Inches(0.6), Inches(1.02), Inches(12.2), Inches(0.35),
        "Several of these controls changed our own conclusions — that is the point.",
        14, MUT)
textbox(s, Inches(0.7), Inches(1.60), Inches(5.9), Inches(4.5),
        "WE CAN SAY\n\n"
        "•  Morphology adds +0.143 over chemical\n    structure for mechanism, 4 of 4 models\n\n"
        "•  It adds nothing for the five toxicity\n    endpoints — chemistry does that\n\n"
        "•  Fusion does not help — we independently\n    replicate the main analysis\n\n"
        "•  The modalities agree weakly but really\n    (adjusted RV 0.016, Mantel p = 0.0001)\n\n"
        "•  Dropping to 2,000 HVGs costs 16× what\n    its leakage was worth",
        14, INK)
textbox(s, Inches(6.9), Inches(1.60), Inches(5.9), Inches(4.5),
        "WE MUST NOT SAY\n\n"
        "•  Any plain RV number — it scores random\n    noise as high as our real data\n\n"
        "•  \"r = 0.903\" — its own null reaches 0.902.\n    Quote the p-value (0.017)\n\n"
        "•  \"Four joint components\" — the null also\n    reached four, p = 0.095\n\n"
        "•  Any gene list — 0 of 41,780 survive\n\n"
        "•  Anything about A549 — its morphology\n    is contaminated (max 1.5e19)",
        14, RED)
textbox(s, Inches(0.7), Inches(6.35), Inches(12.1), Inches(0.9),
        "Every number is shown next to the same number computed on data where the "
        "answer is known to be nothing —\nand several of those controls changed our "
        "conclusions.", 15, ACC, bold=True, align=PP_ALIGN.CENTER)
s.notes_slide.notes_text_frame.text = (
    "CLOSING LINE:\n"
    "\"If you remember one thing: every number we show is printed next to the same "
    "number computed on data where the answer is known to be nothing — and several of "
    "those controls changed our conclusions. At 119 compounds that is worth more than "
    "another decimal place of accuracy.\"\n\n"
    "THE CONTROL THAT KILLED OUR OWN FINDING — say this if you want credibility:\n"
    "AJIVE reported four shared components between morphology and expression. We then "
    "destroyed the compound pairing and it still found four. p = 0.095. So we do not "
    "report it. Volunteering that is worth more than any result on the previous slides.")

prs.save(OUT)
print(f"{OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
